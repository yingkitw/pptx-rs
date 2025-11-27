#!/usr/bin/env python3
"""
Generate a reference PPTX file using python-pptx for comparison with ppt-rs output.
This helps ensure alignment between the two implementations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

def create_reference_presentation():
    """Create a reference presentation matching ppt-rs capabilities."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Set metadata
    prs.core_properties.title = "Alignment Test Presentation"
    prs.core_properties.author = "ppt-rs Team"
    prs.core_properties.subject = "Testing ppt-rs alignment with python-pptx"
    prs.core_properties.keywords = "pptx, rust, python-pptx, alignment"
    prs.core_properties.comments = "This presentation tests alignment between ppt-rs and python-pptx"
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Alignment Test Presentation"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "ppt-rs vs python-pptx Compatibility"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(102, 102, 102)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Content with Shapes
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Shapes and Formatting"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(44)
    title_para2.font.bold = True
    title_para2.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add rectangle shape
    rect = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.8),
        Inches(2.5), Inches(1.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 0, 0)
    rect.line.color.rgb = RGBColor(0, 0, 0)
    rect.line.width = Pt(2)
    
    # Add text to rectangle
    rect_text = rect.text_frame
    rect_text.text = "Rectangle"
    rect_text.paragraphs[0].font.size = Pt(18)
    rect_text.paragraphs[0].font.bold = True
    rect_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    rect_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add circle shape
    circle = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3.5), Inches(1.8),
        Inches(2.5), Inches(1.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0, 255, 0)
    circle.line.color.rgb = RGBColor(0, 0, 0)
    circle.line.width = Pt(2)
    
    # Add text to circle
    circle_text = circle.text_frame
    circle_text.text = "Circle"
    circle_text.paragraphs[0].font.size = Pt(18)
    circle_text.paragraphs[0].font.bold = True
    circle_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    circle_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add diamond shape
    diamond = slide2.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(6.5), Inches(1.8),
        Inches(2.5), Inches(1.5)
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = RGBColor(0, 0, 255)
    diamond.line.color.rgb = RGBColor(0, 0, 0)
    diamond.line.width = Pt(2)
    
    # Add text to diamond
    diamond_text = diamond.text_frame
    diamond_text.text = "Diamond"
    diamond_text.paragraphs[0].font.size = Pt(18)
    diamond_text.paragraphs[0].font.bold = True
    diamond_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    diamond_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_dir = "examples/output"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "reference_python_pptx.pptx")
    prs.save(output_path)
    print(f"✓ Reference presentation created: {output_path}")
    print(f"  - Slides: {len(prs.slides)}")
    print(f"  - Title: {prs.core_properties.title}")
    print(f"  - Author: {prs.core_properties.author}")
    
    return output_path

if __name__ == "__main__":
    try:
        create_reference_presentation()
    except ImportError as e:
        print(f"Error: python-pptx not installed. Install with: pip install python-pptx")
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)

